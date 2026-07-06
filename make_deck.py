"""D9X-15 / D9X-25 — Generate the stakeholder deck: day_9x_stakeholder_deck.pptx.

A detailed, presentation-ready PPT for Vijay to walk stakeholders (or an Anthropic
interviewer) through the prototype. Numbers are pulled LIVE from the model
(breakeven.simulate / team_impact / charter) so the deck always matches the build.

v2.0 adds the knowledge-miner arc: the second money shot (gap → SME interview →
attributed Skill), "Teachers of the sprint" (attribution flips hoarding into
teaching), the one-team impact slide, and a refreshed delivery-evidence slide.

Run:  python make_deck.py
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import config
import charter
import breakeven as be

# ── palette ───────────────────────────────────────────────────────────────────
INK = RGBColor(0x0F, 0x17, 0x2A)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
SKY = RGBColor(0x0E, 0xA5, 0xE9)
GREEN = RGBColor(0x10, 0xB9, 0x81)
SLATE = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W, H = Inches(13.333), Inches(7.5)


def _live_numbers():
    out = {}
    for stars in ("Turnaround", "Start-up", "Sustaining success"):
        e = be.simulate(90, stars)
        d = be.ledger(e)
        out[stars] = {"breakeven": be.breakeven_day(d),
                      "net": float(d["cum_net"].iloc[-1]),
                      "dip": float(d["cum_net"].min())}
    # v2.0 team-impact: analyst-hours Kai returns to the humans over a turnaround run
    impact = be.team_impact(be.simulate(90, "Turnaround"))
    out["hours_returned"] = round(sum(r["Hours returned"] for r in impact), 0)
    return out


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def _run(para, text, size, color, bold=False, italic=False):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Segoe UI"
    return r


def _title_slide(prs, nums):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, INK)
    tf = _box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.2))
    _run(tf.paragraphs[0], "Onboarding Your First AI Teammate", 40, WHITE, bold=True)
    p = tf.add_paragraph()
    _run(p, "day_9x — an AI agent's first 90 days on a credit-risk data-ops team",
         20, RGBColor(0xC7, 0xD2, 0xFE))
    p = tf.add_paragraph()
    _run(p, "v2.0 · the knowledge miner", 13, GREEN, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(18)
    _run(p, "Modelled on Watkins' The First 90 Days · a 6-sprint, eval-gated onboarding · "
            "runs $0 / keyless", 14, SLATE)
    p = tf.add_paragraph()
    p.space_before = Pt(30)
    _run(p, f"Breakeven in sprints, not months — reached day "
            f"{nums['Turnaround']['breakeven']} in a turnaround.", 15, GREEN, italic=True)


def _section(prs, kicker, title, color=INDIGO):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, color)
    tf = _box(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(2))
    _run(tf.paragraphs[0], kicker.upper(), 15, RGBColor(0xE0, 0xE7, 0xFF), bold=True)
    p = tf.add_paragraph()
    _run(p, title, 34, WHITE, bold=True)
    return s


def _content(prs, title, bullets, footer=None):
    """bullets: list of (text, level, color|None)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, WHITE)
    # title bar
    bar = s.shapes.add_shape(1, 0, 0, W, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = INK; bar.line.fill.background()
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.9); tf.margin_top = Inches(0.28)
    _run(tf.paragraphs[0], title, 26, WHITE, bold=True)
    # body
    body = _box(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.4))
    first = True
    for text, level, color in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(8)
        bullet = "•  " if level == 0 else "–  "
        _run(p, bullet + text, 18 - 2 * level, color or INK, bold=(level == 0 and color is None))
    if footer:
        ft = _box(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4))
        _run(ft.paragraphs[0], footer, 11, SLATE, italic=True)
    return s


def _metric_slide(prs, title, cards, footer=None):
    """cards: list of (big, label, color)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, LIGHT)
    tf = _box(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(1))
    _run(tf.paragraphs[0], title, 26, INK, bold=True)
    n = len(cards)
    gap = Inches(0.35)
    cw = (W - Inches(1.8) - gap * (n - 1)) / n
    x = Inches(0.9)
    for big, label, color in cards:
        card = s.shapes.add_shape(1, x, Inches(2.1), cw, Inches(2.6))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.color.rgb = color
        card.line.width = Pt(2)
        ctf = card.text_frame; ctf.word_wrap = True
        ctf.margin_top = Inches(0.5)
        ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(ctf.paragraphs[0], big, 40, color, bold=True)
        p = ctf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        _run(p, label, 14, SLATE)
        x += cw + gap
    if footer:
        ft = _box(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.4))
        _run(ft.paragraphs[0], footer, 15, INK)
    return s


def build():
    nums = _live_numbers()
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # 1 · title
    _title_slide(prs, nums)

    # 2 · the reframe
    _content(prs, "The reframe", [
        ("An AI agent is the fastest-learning new hire you'll ever get — it reads a million "
         "pages before Monday.", 0, None),
        ("And it shows up with zero context about YOUR business: your process, your naming, "
         "your unwritten rules.", 0, None),
        ("Talent was never the gap. Context is.", 0, INDIGO),
        ("So the question isn't 'which model' — it's 'are we a team that knows how to onboard?'",
         0, GREEN),
    ], footer="This is the 'trapped enterprise knowledge' thesis, arriving through the side door.")

    # 3 · why it matters (the numbers)
    _content(prs, "Why onboarding — not the model — decides the outcome", [
        ("Onboarded like a junior teammate (scoped tasks, a reviewing buddy, trust earned "
         "task by task): one vendor's merge rate rose 34% → 67% over 18 months.", 0, GREEN),
        ("Thrown in the deep end: an independent trial completed 3 of 20 tasks — the failure "
         "cause was silent, unpredictable behaviour, not raw error rate.", 0, RGBColor(0xEF,0x44,0x44)),
        ("Watkins: a mid-level human hire takes ~6.2 months to break even; senior outside hires "
         "fail 40–50% of the time. The AI teammate is the ultimate OUTSIDE hire — highest risk, "
         "zero informal network — so it needs the MOST structured onboarding, not the least.", 0, None),
    ], footer="Sources tracked in the POV EVIDENCE file; vendor figures flagged as vendor-reported.")

    # 4 · section — the demo
    _section(prs, "The prototype", "day_9x: five tabs = the first 90 days", INDIGO)

    # 5 · the five tabs
    _content(prs, "The five tabs map to Watkins, mechanism by mechanism", [
        ("Charter — the Five Conversations + a STARS situation selector; hiring manager sets "
         "duration (default 90 days = 6 sprints, 10 checkpoints).", 0, INDIGO),
        ("Shadow mode — RAG over the team's runbooks with citations + an abstain guardrail; the "
         "cultural-learning moment.", 0, INDIGO),
        ("Early wins — scoped 4–8h tasks behind approval gates; helpful-abstain escalation.", 0, SKY),
        ("Probation review — eval gates per task class; trust earned, not assumed.", 0, SKY),
        ("Breakeven + Retro — value consumed vs created; day-9x stakeholder reports.", 0, GREEN),
    ], footer="Runs $0 / keyless (public-Space safe); optional Claude Haiku enriches shadow answers.")

    # 6 · the trapped-knowledge moment
    _content(prs, "The money shot: unlocking trapped knowledge in one click", [
        ("Ask Kai: 'Can I rerun the CreditMart load on Thursday morning?'", 0, None),
        ("Before: it can't — the 'Thursday rule' is in nobody's document, only in an SME's head. "
         "It answers from the nearest runbook and misses it.", 1, SLATE),
        ("An SME coaches the unwritten rule → it becomes a versioned, cited Skill in the corpus.", 0, None),
        ("After: Kai answers correctly, citing unwritten_rules.md.", 1, GREEN),
        ("The probation-review evals prove it objectively: two culture cases are RED until coached, "
         "then GREEN. Evals, not vibes.", 0, INDIGO),
    ], footer="Trapped knowledge → reusable Skill. This is the thesis, made demoable.")

    # 6b · the SECOND money shot — the knowledge miner (v2.0)
    _content(prs, "The second money shot: Kai mines what the corpus can't answer", [
        ("When Kai can't answer, the miss isn't lost — it's logged as a structured "
         "knowledge gap (a mining lead), never a silent hallucination.", 0, None),
        ("Kai turns the gap into SME interview questions; the SME answers; Kai writes it "
         "up as a versioned, ATTRIBUTED Skill and reindexes — the question it abstained "
         "on now answers, with a citation.", 0, INDIGO),
        ("This is agent-led externalization: the 80% that was never written down gets "
         "captured by interviewing the human who has it (research: 94.9% recall in "
         "simulation), and the SME is the credited teacher — never an extraction target.",
         0, None),
        ("Almost nobody demos this. It is exactly the tacit-knowledge loop the POV names "
         "as the prototype opportunity.", 0, GREEN),
    ], footer="Gap register → SME interview → cited, versioned Skill. Trapped knowledge, "
              "unlocked with consent and attribution.")

    # 6c · teachers of the sprint — attribution as the incentive (v2.0)
    _content(prs, "Teachers of the sprint — why attribution is the whole game", [
        ("35% of knowledge workers deliberately hoard expertise to stay indispensable in "
         "the AI era; AI awareness measurably increases knowledge-hiding.", 0, RGBColor(0xEF,0x44,0x44)),
        ("So a program pitched as 'we'll capture the SMEs' knowledge' poisons its own "
         "corpus. The fix is design, not sentiment: SMEs as named, credited teachers.", 0, None),
        ("day_9x makes it visible — a 'Teachers of the sprint' board counts Skills coached "
         "and citations served per SME: 'your knowledge answered N questions this sprint.'", 0, INDIGO),
        ("Where AI is embedded well, 48% of workers feel energized vs 19% without. "
         "Attribution is corpus-quality engineering.", 0, GREEN),
    ], footer="The 'one team' framing isn't a slogan — it's what keeps the knowledge flowing in.")

    # 7 · breakeven metrics (live)
    _metric_slide(prs, "The breakeven point — in sprints, not months", [
        (f"Day {nums['Turnaround']['breakeven']}", "Turnaround breakeven\n(short shadow, fast wins)", GREEN),
        (f"Day {nums['Sustaining success']['breakeven']}", "Sustaining-success breakeven", SKY),
        (f"Day {nums['Start-up']['breakeven']}", "Start-up breakeven\n(deep learning first)", INDIGO),
    ], footer="The STARS situation visibly changes WHEN the AI teammate pays back its onboarding "
              f"cost. Net at day 90 (turnaround): +{nums['Turnaround']['net']:.0f} analyst-hours. "
              "The curve dips through shadow (drafts create no realized value — the human still "
              "owns the work), then climbs as scoped work is approved.")

    # 8 · trust model
    _content(prs, "Trust is earned per task class — and revocable", [
        ("SHADOW — drafts only, 100% reviewed, zero autonomy.", 0, INDIGO),
        ("GATED — may act, but every output needs the buddy's approval; unlocked when the class "
         "passes its eval gate.", 0, SKY),
        ("AUTONOMOUS — acts within the charter; ~20% audit sample; offered only with a green gate "
         "AND enough verified outputs.", 0, GREEN),
        ("A human junior earns social trust; an agent earns statistical trust. The probation "
         "review can expand, hold, or reduce — the same 90-day decision you'd make for a person.",
         0, None),
    ], footer="What stays human, in writing: regulatory sign-off, stakeholder calls, ambiguous "
              "judgement, accountability.")

    # 8b · one team — what each human gets back (v2.0)
    _content(prs, "One team — what each human gets back", [
        ("In this story the AI is the junior; the humans are the mentors, reviewers and "
         "deciders. Every hour Kai returns is redeployed toward the review-and-judgement "
         "work that was always understaffed.", 0, None),
        (f"Over a modelled turnaround run, Kai returns ~{nums['hours_returned']:.0f} "
         "analyst-hours to named teammates — Priya (triage), Arun (recon), Sofia (DQ) — who "
         "shift from doing the routine work to reviewing and directing it.", 0, INDIGO),
        ("The 'stays human' list is written down, not implied: regulatory sign-off, "
         "stakeholder calls, ambiguous judgement, accountability.", 0, None),
        ("Resist FTE-equivalence math in year one — measure task-class throughput, quality "
         "and cycle time. The workforce story is redeployment, not reduction.", 0, GREEN),
    ], footer="G3 answered in the product, not just the pitch: the 52%-job-fear stat met with a "
              "per-person impact view and a written human-value-preservation list.")

    # 9 · delivery evidence
    _content(prs, "Delivery evidence — built the way the team works", [
        ("Jira board (project D9X): 25 stories across 4 sprints, story-first, AI-actual worklogs "
         "— importable via JIRA_IMPORT.csv. v0.1 frozen (git tag), v2.0 shipped.", 0, None),
        ("v2.0 realism, not just narrative: a freshness rule excludes superseded runbooks; an "
         "ACL refusal keeps salary/HR out of scope; a retrieval-eval scorecard grades whether the "
         "RIGHT knowledge surfaces (most RAG teams run no retrieval evals at all).", 0, SKY),
        ("Tests: 23/23 green. Eval harness: 2 culture gates red uncoached → all green coached "
         "(the freshness + ACL cases pass green throughout — they are realism, not the teaching "
         "moment).", 0, GREEN),
        ("Built with a partner-agent team (Opus subagents for docs + an independent review pass, "
         "D9X-26); token usage tracked per sprint. GitHub + HF Space deploy ready; push/deploy "
         "pending a PAT / HF token supplied at the time (deliberately not stored).", 0, None),
    ], footer="Governance rails: synthetic data only; no key on a public Space; Watkins paraphrased "
              "with attribution; an independent review documents outcomes in REVIEW_v2.md.")

    # 10 · close
    s = _section(prs, "The close", "Not a job-loss story — a promotion story", GREEN)
    tf = _box(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(2))
    _run(tf.paragraphs[0],
         "In this story the AI is the junior. Your people are the mentors, the reviewers, the ones "
         "who decide. We're applying the management playbook they already trust — to a new kind of "
         "teammate.", 18, WHITE)
    p = tf.add_paragraph(); p.space_before = Pt(14)
    _run(p, config.WATKINS_CREDIT, 11, RGBColor(0xD1, 0xFA, 0xE5), italic=True)

    out = config.ROOT / "day_9x_stakeholder_deck.pptx"
    prs.save(out)
    return out, nums


if __name__ == "__main__":
    path, nums = build()
    print(f"Deck written: {path}")
    print(f"  slides: 13 | live breakeven (turnaround): day {nums['Turnaround']['breakeven']}, "
          f"net +{nums['Turnaround']['net']:.0f}h | hours returned to humans: "
          f"~{nums['hours_returned']:.0f}h")
